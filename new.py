from flask import Flask, render_template, request, session, redirect, flash, send_file
from flask_sqlalchemy import SQLAlchemy
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.asymmetric import dh
from cryptography.hazmat.backends import default_backend
from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
from cryptography.hazmat.primitives import padding
import os
import docx
from pptx import Presentation
import json
import io
import base64

# Load configuration from JSON
with open('config.json', 'r') as c:
    params = json.load(c)["params"]

local_server = True
app = Flask(__name__)
app.secret_key = 'super-secret-key'

# Configure database URI
if local_server:
    app.config['SQLALCHEMY_DATABASE_URI'] = params['local_uri']
else:
    app.config['SQLALCHEMY_DATABASE_URI'] = params['proud_uri']

db = SQLAlchemy(app)

# Models based on the schema in db.sql
class Patients(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    patient_id = db.Column(db.Integer, nullable=False)
    patient_name = db.Column(db.String(255), nullable=False)
    department = db.Column(db.String(50), nullable=False)
    comments = db.Column(db.Text)
    file_path = db.Column(db.String(255))  # Store the file path
    upload_date = db.Column(db.TIMESTAMP, server_default=db.func.current_timestamp())
    encrypted_content = db.Column(db.Text)  # Store the encrypted text
    encryption_key = db.Column(db.Text)  # Store the encryption key


# Generate Diffie-Hellman key exchange
def generate_key():
    parameters = dh.generate_parameters(generator=2, key_size=512, backend=default_backend())
    private_key = parameters.generate_private_key()
    public_key = private_key.public_key()
    return private_key, public_key


# Encrypt the document content
def encrypt_text(key, text):
    iv = os.urandom(16)
    cipher = Cipher(algorithms.AES(key), modes.CBC(iv), backend=default_backend())
    encryptor = cipher.encryptor()

    padder = padding.PKCS7(algorithms.AES.block_size).padder()
    padded_data = padder.update(text.encode()) + padder.finalize()
    encrypted = encryptor.update(padded_data) + encryptor.finalize()

    return base64.b64encode(iv + encrypted).decode('utf-8')


# Flask Routes
@app.route("/")
def hello():
    return render_template('index.html', params=params)


@app.route("/index")
def home():
    return render_template('dashboard.html', params=params)


@app.route("/search", methods=['GET', 'POST'])
def search():
    if request.method == 'POST':
        name = request.form.get('search')
        patient = Patients.query.filter_by(patient_name=name).first()
        if patient:
            flash("Patient Found.", "primary")
        else:
            flash("Patient Not Found.", "danger")
    return render_template('search.html', params=params)


@app.route("/details", methods=['GET', 'POST'])
def details():
    if 'user' in session and session['user'] == params['user']:
        logs = ActionLogs.query.all()
        return render_template('details.html', params=params, logs=logs)


@app.route("/aboutus")
def aboutus():
    return render_template('aboutus.html', params=params)


@app.route("/insert", methods=['GET', 'POST'])
def insert():
    if request.method == 'POST':
        patient_id = request.form.get('patient_id')
        patient_name = request.form.get('patient_name')
        department = request.form.get('department')
        comments = request.form.get('comments')
        document_file = request.files['document']

        if not patient_id or not patient_name:
            flash("Patient ID and Name are required!", "danger")
            return render_template('insert.html', params=params)

        # Save the uploaded document file
        if document_file:
            filename = document_file.filename
            file_path = os.path.join("uploads", filename)  # Specify the uploads directory
            document_file.save(file_path)
        else:
            flash("No file uploaded!", "danger")
            return render_template('insert.html', params=params)

        # Process and encrypt the document content
        document_text = ''
        if filename.endswith('.docx'):
            doc = docx.Document(file_path)
            document_text = '\n'.join([para.text for para in doc.paragraphs])
        elif filename.endswith('.pptx'):
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        document_text += shape.text + '\n'
        else:
            flash("Unsupported file type!", "danger")
            return render_template('insert.html', params=params)

        # Generate a random encryption key
        encryption_key = os.urandom(32)  # Generate a random key for AES
        encrypted_text = encrypt_text(encryption_key, document_text)

        # Include the encryption key in the encrypted content (base64 encode it for storage)
        encoded_key = base64.b64encode(encryption_key).decode('utf-8')
        final_encrypted_content = f"{encoded_key}:{encrypted_text}"

        new_patient = Patients(
            patient_id=patient_id,
            patient_name=patient_name,
            department=department,
            comments=comments,
            file_path=file_path,  # Store the file path
            encrypted_content=final_encrypted_content  # Store the key and encrypted content together
        )
        db.session.add(new_patient)
        db.session.commit()

        flash("Patient and Encrypted Document Added Successfully", "success")
        return render_template('insert.html', params=params)

    return render_template('insert.html', params=params)


@app.route("/receive", methods=['GET', 'POST'])
def receive():
    data_available = False
    patient = None
    if request.method == 'POST':
        patient_id = request.form.get('patient_id')
        patient_name = request.form.get('patient_name')
        department = request.form.get('department')
        feedback = request.form.get('feedback')

        # Logic to check if the selected department has sent any data
        patient = Patients.query.filter_by(patient_id=patient_id, patient_name=patient_name).first()
        if patient and patient.department == department:  # Check department as well
            data_available = True
            flash("Data found for the selected patient.", "success")
        else:
            flash("No data available for the selected patient from this department.", "danger")

        return render_template('receive.html', params=params, data_available=data_available, received_data=patient)

    return render_template('receive.html', params=params, data_available=data_available)


@app.route("/logout")
def logout():
    session.pop('user')
    flash("You are logged out", "primary")
    return redirect('/login')


@app.route("/login", methods=['GET', 'POST'])
def login():
    if 'user' in session and session['user'] == params['user']:
        return render_template('dashboard.html', params=params)

    if request.method == 'POST':
        username = request.form.get('uname')
        password = request.form.get('password')
        if username == params['user'] and password == params['password']:
            session['user'] = username
            flash("You are logged in", "primary")
            return redirect('/index')
        else:
            flash("Invalid credentials", "danger")
    return render_template('login.html', params=params)


@app.route("/download/<int:patient_id>")
def download_encrypted(patient_id):
    patient = Patients.query.filter_by(patient_id=patient_id).first()
    if patient:
        # Return the original uploaded file
        if patient.file_path:
            return send_file(patient.file_path, as_attachment=True)
        else:
            flash("No file uploaded for this patient.", "danger")
            return redirect('/')
    else:
        flash("Patient not found!", "danger")
        return redirect('/')


def decrypt_text(key, encrypted_text):
    encrypted_data = base64.b64decode(encrypted_text.encode('utf-8'))
    iv = encrypted_data[:16]
    encrypted_data = encrypted_data[16:]

    cipher = Cipher(algorithms.AES(key), modes.CBC(iv), backend=default_backend())
    decryptor = cipher.decryptor()

    decrypted_padded = decryptor.update(encrypted_data) + decryptor.finalize()

    unpadder = padding.PKCS7(algorithms.AES.block_size).unpadder()
    decrypted = unpadder.update(decrypted_padded) + unpadder.finalize()

    return decrypted.decode('utf-8')


# Start the application
if __name__ == "__main__":
    os.makedirs("uploads", exist_ok=True)
    app.run(debug=True)
